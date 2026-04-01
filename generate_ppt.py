from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

prs = Presentation()

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_bullet_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)

# Slide 1: Title
add_title_slide(prs, "Lumina AI Synthesizer", "Intelligent Document and Media Summarization\nSimple, Fast, and Powerful")

# Slide 2: What is Lumina?
add_bullet_slide(prs, "The Challenge & Solution", [
    "Challenge: We are flooded with long texts, large documents, images, and videos. Reviewing them all manually takes hours.",
    "Solution: Lumina! It instantly reads, watches, or analyzes your files to distil the most important information.",
    "Everything happens automatically in one place, using a beautiful and intuitive interface designed for ease of use."
])

# Slide 3: Key Features
add_bullet_slide(prs, "What Lumina Can Do", [
    "1. Text Summarization: Paste entire articles or essays and get instant summaries.",
    "2. Document Understanding: Upload massive PDFs or Word Docs for rapid insights.",
    "3. Media Analysis: Drop in images, GIFs, or Videos, and Lumina describes and understands them.",
    "4. Clean Export: Download your insights easily to PDF, Text, or Word formats.",
    "5. Memory Timeline: Retrieve previous analysis from your history securely and instantly."
])

# Slide 4: Workflow Diagram
slide_layout = prs.slide_layouts[5] # Title only
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "How It Works (The Workflow)"

# Blocks
lefts = [Inches(1), Inches(4), Inches(7)]
top = Inches(3)
width = Inches(2)
height = Inches(1.2)
texts = ["1. Upload Data\n(Drop Texts, PDFs,\nImages, or Videos)", "2. Auto Engine\n(Intelligently analyzes\nthe input)", "3. Receive Output\n(Read, Export, or\nSave to History)"]

from pptx.enum.text import PP_ALIGN

for i in range(3):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lefts[i], top, width, height)
    shape.fill.solid()
    # Nice warm brand color
    shape.fill.fore_color.rgb = RGBColor(212, 162, 74)
    shape.line.color.rgb = RGBColor(180, 130, 40)
    
    tf = shape.text_frame
    tf.text = texts[i]
    for p in tf.paragraphs:
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(26, 20, 16)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

# Arrows
for i in range(2):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, lefts[i] + width + Inches(0.2), top + Inches(0.4), Inches(0.6), Inches(0.4))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(160, 160, 160)

# Slide 5: The Smart Engine Under the Hood
add_bullet_slide(prs, "The Smart Engine Under the Hood", [
    "Lumina has two 'brains' running seamlessly without any complex settings:",
    "Brain A (The Reader): Exceptionally fast at reading words from text and multiple documents rapidly to give you the gist.",
    "Brain B (The Watcher): Specially designed to 'look' at rich media (pictures and videos) and describe what's happening.",
    "Lumina is smart enough to choose the right brain automatically based on what files you upload."
])

# Slide 6: Conclusion
add_bullet_slide(prs, "Why Choose Lumina?", [
    "A completely hassle-free experience with absolutely zero jargon—it just works out of the box.",
    "Reclaims your time from hours of tedious reading and manual review.",
    "Modern, premium design that feels alive and keeps you focused.",
    "Always remembers your past work so you don't lose that stroke of genius."
])

prs.save("Lumina_Project_Presentation.pptx")
print("Successfully generated Lumina_Project_Presentation.pptx")
